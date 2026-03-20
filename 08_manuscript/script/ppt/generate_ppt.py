#!/usr/bin/env python3
"""
Generate a clean PowerPoint deck from a YAML outline.

Designed for a simple workflow:
  1) Edit 08_manuscript/yaml/ppt.yaml
  2) Run this script
  3) Get 08_manuscript/out/<name>.pptx

Recommended installation inside a micromamba environment:
  python -m pip install python-pptx pyyaml pillow cairosvg

Usage:
  python generate_ppt.py /path/to/ppt.yaml
  python generate_ppt.py /path/to/ppt.yaml /path/to/output.pptx

YAML expectation (flexible):
Top-level keys are sections, each mapping to a list of slides.
Each slide can contain:
  - id: 1
  - type: title | intro | result | model | summary | section
  - title: ...
  - figure: /path/to/file.png OR [list, of, files]
  - message: optional string or list of strings
  - notes: optional string or list

This script is intentionally conservative and robust:
  - missing figures become placeholders instead of hard failures
  - SVG figures are rasterized to PNG when cairosvg is available
  - list messages are rendered as bullets or separate lines
"""
from __future__ import annotations

import argparse
import hashlib
import os
import re
import shutil
import sys
import tempfile
from pathlib import Path
from typing import Any, Dict, Iterable, List, Sequence, Tuple

try:
    import yaml
except ImportError as e:
    raise SystemExit("Missing dependency: PyYAML. Install with `python -m pip install pyyaml`.") from e

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError as e:
    raise SystemExit("Missing dependency: python-pptx. Install with `python -m pip install python-pptx`.") from e

try:
    from PIL import Image
except ImportError as e:
    raise SystemExit("Missing dependency: Pillow. Install with `python -m pip install pillow`.") from e

try:
    import cairosvg  # type: ignore
except Exception:
    cairosvg = None

# ----------------------------
# Theme
# ----------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_X = Inches(0.55)
TOP_BAND_Y = Inches(0.35)
TOP_BAND_H = Inches(0.75)
CONTENT_Y = Inches(1.3)
BOTTOM_BOX_H = Inches(0.85)
BOTTOM_GAP = Inches(0.18)
FOOTER_Y = SLIDE_H - Inches(0.22)
ACCENT = RGBColor(79, 129, 189)      # muted Nature-like blue
ACCENT_LIGHT = RGBColor(232, 239, 247)
INK = RGBColor(17, 24, 39)           # gray-900
MUTED = RGBColor(107, 114, 128)      # gray-500
BORDER = RGBColor(209, 213, 219)     # gray-300
BG = RGBColor(255, 255, 255)
SOFT = RGBColor(249, 250, 251)       # gray-50

FONT_CJK = "Arial Unicode MS"
FONT_LATIN = "Aptos"
DEFAULT_FONT = FONT_CJK

# ----------------------------
# Utilities
def emu(value):
    """Convert small numeric layout values to EMU integers for python-pptx.

    python-pptx text frame margins must be integral EMU values, not floats.
    Convention in this script: bare ints/floats are interpreted as inches.
    """
    if value is None:
        return 0
    if isinstance(value, int):
        return value
    if isinstance(value, float):
        return Inches(value)
    return value
# ----------------------------
def as_list(value: Any) -> List[Any]:
    if value is None:
        return []
    if isinstance(value, list):
        return value
    return [value]


def normalize_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, list):
        return "\n".join(str(v) for v in value if v is not None)
    return str(value)


def safe_filename(text: str) -> str:
    text = re.sub(r"[^\w\-\.]+", "_", text.strip())
    return text or "ppt_from_yaml_auto"


def infer_output_path(input_yaml: Path) -> Path:
    out_dir = input_yaml.parent.parent / "out"
    out_dir.mkdir(parents=True, exist_ok=True)
    return out_dir / "ppt_from_yaml_auto.pptx"


def iter_sections(data: Dict[str, Any]) -> Iterable[Tuple[str, List[Dict[str, Any]]]]:
    for key, value in data.items():
        if isinstance(value, list):
            slides = [item for item in value if isinstance(item, dict)]
            if slides:
                yield str(key), slides


def title_case_section(name: str) -> str:
    return str(name).strip()


def fit_font_size(title: str, max_size: int = 26, min_size: int = 18) -> int:
    n = len(title)
    if n <= 20:
        return max_size
    if n <= 32:
        return 24
    if n <= 44:
        return 22
    if n <= 60:
        return 20
    return min_size


def add_textbox(slide, text: str, x, y, w, h, *, font_size=20, bold=False,
                color=INK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
                margin_left=0.04, margin_right=0.04, margin_top=0.02, margin_bottom=0.02,
                font_name=DEFAULT_FONT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = emu(margin_left)
    tf.margin_right = emu(margin_right)
    tf.margin_top = emu(margin_top)
    tf.margin_bottom = emu(margin_bottom)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    return box


def add_bullets(slide, lines: Sequence[str], x, y, w, h, *, font_size=18,
                color=INK, bullet=True, font_name=DEFAULT_FONT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = emu(0.04)
    tf.margin_right = emu(0.04)
    tf.margin_top = emu(0.02)
    tf.margin_bottom = emu(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate([str(v) for v in lines if str(v).strip()]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.bullet = bullet
        font = p.font
        font.name = font_name
        font.size = Pt(font_size)
        font.bold = False
        font.color.rgb = color
    return box


def add_top_band(slide, title: str, slide_id: Any = None):
    # subtle top rule
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, TOP_BAND_Y + TOP_BAND_H - Inches(0.04), SLIDE_W, Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # title
    add_textbox(
        slide,
        title,
        MARGIN_X,
        TOP_BAND_Y,
        SLIDE_W - 2 * MARGIN_X,
        TOP_BAND_H - Inches(0.06),
        font_size=fit_font_size(title),
        bold=True,
        color=INK,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_footer(slide, left_text: str = ""):
    if left_text:
        add_textbox(slide, left_text, MARGIN_X, FOOTER_Y - Inches(0.06), Inches(5), Inches(0.16),
                    font_size=9, color=MUTED, valign=MSO_ANCHOR.MIDDLE)


def is_svg(path: Path) -> bool:
    return path.suffix.lower() == ".svg"


def rasterize_if_needed(path: Path, cache_dir: Path) -> Path:
    if not path.exists():
        return path
    if not is_svg(path):
        return path
    if cairosvg is None:
        return path
    key = hashlib.md5(str(path).encode("utf-8")).hexdigest()[:12]
    out = cache_dir / f"{path.stem}_{key}.png"
    if out.exists() and out.stat().st_mtime >= path.stat().st_mtime:
        return out
    out.parent.mkdir(parents=True, exist_ok=True)
    cairosvg.svg2png(url=str(path), write_to=str(out), output_width=2200)
    return out


def get_image_size(path: Path) -> Tuple[int, int]:
    with Image.open(path) as img:
        return img.size


def contain_fit(box_x, box_y, box_w, box_h, img_w, img_h):
    img_ratio = img_w / img_h
    box_ratio = box_w / box_h
    if img_ratio > box_ratio:
        w = box_w
        h = box_w / img_ratio
        x = box_x
        y = box_y + (box_h - h) / 2
    else:
        h = box_h
        w = box_h * img_ratio
        x = box_x + (box_w - w) / 2
        y = box_y
    return x, y, w, h


def add_placeholder(slide, x, y, w, h, text="Figure placeholder"):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = SOFT
    shp.line.color.rgb = BORDER
    add_textbox(slide, text, x + Inches(0.2), y + Inches(0.2), w - Inches(0.4), h - Inches(0.4),
                font_size=16, color=MUTED, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)




def add_figures(slide, figures: Any, x, y, w, h, cache_dir: Path):
    fig_list = [Path(str(f)).expanduser() for f in as_list(figures) if str(f).strip()]
    valid = []
    missing = []
    for fig in fig_list:
        fig2 = rasterize_if_needed(fig, cache_dir)
        if fig2.exists() and fig2.is_file() and fig2.suffix.lower() in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tif", ".tiff"}:
            valid.append(fig2)
        elif fig.exists() and is_svg(fig) and cairosvg is None:
            missing.append(f"SVG unsupported (install cairosvg): {fig.name}")
        else:
            missing.append(f"Missing: {fig}")

    if not valid:
        msg = "Figure placeholder"
        if missing:
            msg += "\n" + "\n".join(missing[:3])
        add_placeholder(slide, x, y, w, h, msg)
        return

    n = len(valid)
    if n == 1:
        path = valid[0]
        iw, ih = get_image_size(path)
        px, py, pw, ph = contain_fit(x, y, w, h, iw, ih)
        slide.shapes.add_picture(str(path), px, py, pw, ph)
        return

    if n == 2:
        gap = Inches(0.15)
        col_w = (w - gap) / 2
        boxes = [(x, y, col_w, h), (x + col_w + gap, y, col_w, h)]
    else:
        gap = Inches(0.15)
        cols = 2
        rows = (n + 1) // 2
        cell_w = (w - gap) / 2
        cell_h = (h - gap * (rows - 1)) / rows
        boxes = []
        for i in range(n):
            r = i // cols
            c = i % cols
            boxes.append((x + c * (cell_w + gap), y + r * (cell_h + gap), cell_w, cell_h))

    for path, (bx, by, bw, bh) in zip(valid, boxes):
        iw, ih = get_image_size(path)
        px, py, pw, ph = contain_fit(bx, by, bw, bh, iw, ih)
        slide.shapes.add_picture(str(path), px, py, pw, ph)


def render_title_slide(prs: Presentation, slide_data: Dict[str, Any], cache_dir: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Accent bar
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(0.65), Inches(0.18), Inches(1.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    title = normalize_text(slide_data.get("title"))
    add_textbox(slide, title, Inches(0.95), Inches(0.8), Inches(10.6), Inches(1.7),
                font_size=28 if len(title) < 28 else 24, bold=True, color=INK)

    message = slide_data.get("message")
    msg_lines = as_list(message) if isinstance(message, list) else [normalize_text(message)] if normalize_text(message).strip() else []
    if msg_lines:
        add_bullets(slide, msg_lines, Inches(0.98), Inches(2.35), Inches(5.8), Inches(1.4), font_size=18, bullet=False, color=MUTED)

    figs = slide_data.get("figure")
    if figs:
        add_figures(slide, figs, Inches(7.9), Inches(1.15), Inches(4.5), Inches(4.8), cache_dir)

    add_footer(slide, "Auto-generated from ppt.yaml")
    return slide


def render_section_divider(prs: Presentation, section_name: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(2.0), SLIDE_W, Inches(1.8))
    strip.fill.solid()
    strip.fill.fore_color.rgb = SOFT
    strip.line.fill.background()

    add_textbox(slide, title_case_section(section_name), Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.9),
                font_size=28, bold=True, color=INK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.9), Inches(3.3), Inches(2.2), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    return slide


def render_content_slide(prs: Presentation, slide_data: Dict[str, Any], cache_dir: Path, section_name: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    title = normalize_text(slide_data.get("title"))
    slide_id = slide_data.get("id")
    add_top_band(slide, title, slide_id)

    # Body geometry
    body_h = SLIDE_H - CONTENT_Y - Inches(0.45)

    add_figures(slide, slide_data.get("figure"), MARGIN_X, CONTENT_Y, SLIDE_W - 2 * MARGIN_X, body_h, cache_dir)

    add_footer(slide, section_name)
    return slide


def build_presentation(data: Dict[str, Any], output_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    cache_dir = output_path.parent / ".ppt_cache"
    cache_dir.mkdir(parents=True, exist_ok=True)

    sections = list(iter_sections(data))
    if not sections:
        raise ValueError("YAML 顶层没有找到任何 section:list 结构。")

    first_section_name, first_slides = sections[0]
    # If the very first slide is a title slide, use it directly.
    if first_slides and str(first_slides[0].get("type", "")).lower() == "title":
        render_title_slide(prs, first_slides[0], cache_dir)
        first_slides = first_slides[1:]

    # Remaining slides of first section, no divider by default.
    for slide_data in first_slides:
        render_content_slide(prs, slide_data, cache_dir, first_section_name)

    for section_name, slides in sections[1:]:
        render_section_divider(prs, section_name)
        for slide_data in slides:
            if str(slide_data.get("type", "")).lower() == "title":
                render_title_slide(prs, slide_data, cache_dir)
            else:
                render_content_slide(prs, slide_data, cache_dir, section_name)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def main():
    parser = argparse.ArgumentParser(description="Generate a PowerPoint deck from a YAML file.")
    parser.add_argument("yaml_path", help="Path to ppt.yaml")
    parser.add_argument("output_path", nargs="?", help="Optional output .pptx path")
    args = parser.parse_args()

    yaml_path = Path(args.yaml_path).expanduser().resolve()
    if not yaml_path.exists():
        raise SystemExit(f"YAML file not found: {yaml_path}")

    output_path = Path(args.output_path).expanduser().resolve() if args.output_path else infer_output_path(yaml_path)

    with open(yaml_path, "r", encoding="utf-8") as f:
        data = yaml.safe_load(f)
    if not isinstance(data, dict):
        raise SystemExit("Top-level YAML must be a mapping/dictionary.")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    build_presentation(data, output_path)
    print(f"Saved PPT to: {output_path}")
    print("Tip: if SVG figures do not render, install cairosvg in your micromamba environment.")


if __name__ == "__main__":
    main()
os.system(f"open 08_manuscript/out/ppt_from_yaml_auto.pptx")